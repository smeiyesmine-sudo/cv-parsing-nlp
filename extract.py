"""Text extraction from CV files.

Supported formats:
  .pdf              -> PyMuPDF text layer; OCR fallback for scanned PDFs
  .docx             -> python-docx (paragraphs + tables)
  .doc  (legacy)    -> auto-convert to .docx (LibreOffice or MS Word), then parse
  .pptx             -> python-pptx (all shapes + tables on all slides)
  .ppt  (legacy)    -> auto-convert to .pptx, then parse
  .png .jpg .jpeg   -> OCR with Tesseract (French + English)
  .txt .md          -> read directly

OCR requires Tesseract installed on the system:
  Windows: https://github.com/UB-Mannheim/tesseract/wiki
           (tick "French" under Additional language data during install)
"""
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

# Minimum characters for a PDF text layer to be considered "real text";
# below this we assume it's a scanned document and run OCR.
PDF_MIN_TEXT_CHARS = 100

# Default Tesseract install locations on Windows
_WIN_TESSERACT_PATHS = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]


def extract_text(path: str) -> str:
    """Extract raw text from a CV file. Dispatches on extension."""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")
    ext = p.suffix.lower()
    if ext == ".pdf":
        return _from_pdf(p)
    if ext == ".docx":
        return _from_docx(p)
    if ext == ".doc":
        return _from_docx(_convert_legacy(p, "docx"))
    if ext == ".pptx":
        return _from_pptx(p)
    if ext == ".ppt":
        return _from_pptx(_convert_legacy(p, "pptx"))
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"):
        return _ocr_image_file(p)
    if ext in (".txt", ".md"):
        return p.read_text(encoding="utf-8", errors="replace")
    raise ValueError(
        f"Unsupported file type: '{ext}'. "
        "Supported: .pdf .docx .doc .pptx .ppt .png .jpg .jpeg .txt"
    )


# ---------------------------------------------------------------- PDF

def _from_pdf(p: Path) -> str:
    import pymupdf
    doc = pymupdf.open(str(p))
    pages = [page.get_text() for page in doc]
    text = "\n".join(pages)
    if len(text.strip()) >= PDF_MIN_TEXT_CHARS:
        doc.close()
        return text
    # Scanned PDF (no real text layer) -> render pages and OCR them
    ocr_pages = []
    for page in doc:
        pix = page.get_pixmap(dpi=300)
        ocr_pages.append(_ocr_pixmap(pix))
    doc.close()
    return "\n".join(ocr_pages)


def _ocr_pixmap(pix) -> str:
    import io
    from PIL import Image
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    return _ocr_pil_image(img)


# --------------------------------------------------------------- DOCX

def _from_docx(p: Path) -> str:
    import docx
    doc = docx.Document(str(p))
    parts = [para.text for para in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts)


# --------------------------------------------------------------- PPTX

def _from_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
    return "\n".join(parts)


# ------------------------------------------------------- Images / OCR

def _configure_tesseract():
    """On Windows, point pytesseract at the default install path if
    tesseract isn't already on PATH."""
    import pytesseract
    if shutil.which("tesseract"):
        return
    if sys.platform.startswith("win"):
        for cand in _WIN_TESSERACT_PATHS:
            if os.path.exists(cand):
                pytesseract.pytesseract.tesseract_cmd = cand
                return
    raise RuntimeError(
        "Tesseract OCR is not installed (needed for images and scanned PDFs).\n"
        "Windows installer: https://github.com/UB-Mannheim/tesseract/wiki\n"
        "During installation, tick 'French' under additional language data."
    )


def _ocr_langs() -> str:
    """Use fra+eng if French data is installed, else eng."""
    import pytesseract
    try:
        available = set(pytesseract.get_languages(config=""))
    except Exception:
        available = set()
    return "fra+eng" if "fra" in available else "eng"


def _ocr_pil_image(img) -> str:
    import pytesseract
    _configure_tesseract()
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    return pytesseract.image_to_string(img, lang=_ocr_langs())


def _ocr_image_file(p: Path) -> str:
    from PIL import Image
    with Image.open(p) as img:
        return _ocr_pil_image(img)


# --------------------------------------------- Legacy .doc / .ppt

def _convert_legacy(p: Path, target: str) -> Path:
    """Convert legacy .doc/.ppt to .docx/.pptx.

    Tries, in order:
      1. LibreOffice (soffice) if installed
      2. MS Word / PowerPoint COM automation (Windows with Office installed)
    The converted file is cached next to a temp dir and reused.
    """
    out_dir = Path(tempfile.gettempdir()) / "cv_parser_converted"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / (p.stem + "." + target)
    if out_path.exists() and out_path.stat().st_mtime >= p.stat().st_mtime:
        return out_path  # already converted

    # --- try LibreOffice
    soffice = (shutil.which("soffice")
               or _first_existing([
                   r"C:\Program Files\LibreOffice\program\soffice.exe",
                   r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
               ]))
    if soffice:
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", target,
                 "--outdir", str(out_dir), str(p)],
                check=True, capture_output=True, timeout=120,
            )
            if out_path.exists():
                return out_path
        except Exception:
            pass

    # --- try MS Office COM automation (Windows only)
    if sys.platform.startswith("win"):
        converted = _convert_with_ms_office(p, out_path, target)
        if converted:
            return out_path

    raise RuntimeError(
        f"Cannot convert legacy file '{p.name}'. Either:\n"
        f"  - install LibreOffice (https://www.libreoffice.org), or\n"
        f"  - open the file in Word/PowerPoint and save it as .{target}"
    )


def _convert_with_ms_office(p: Path, out_path: Path, target: str) -> bool:
    try:
        import win32com.client  # requires: pip install pywin32
    except ImportError:
        return False
    try:
        if target == "docx":
            app = win32com.client.Dispatch("Word.Application")
            app.Visible = False
            doc = app.Documents.Open(str(p.resolve()))
            doc.SaveAs2(str(out_path.resolve()), FileFormat=16)  # wdFormatDocumentDefault
            doc.Close(False)
            app.Quit()
        else:  # pptx
            app = win32com.client.Dispatch("PowerPoint.Application")
            pres = app.Presentations.Open(str(p.resolve()), WithWindow=False)
            pres.SaveAs(str(out_path.resolve()), 24)  # ppSaveAsOpenXMLPresentation
            pres.Close()
            app.Quit()
        return out_path.exists()
    except Exception:
        return False


def _first_existing(paths):
    for c in paths:
        if os.path.exists(c):
            return c
    return None
